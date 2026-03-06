"""
MoMSME Document Formatter – Full Feature Edition v5.2

NEW in v5.2:
- Color theme selector: MoMSME, KPMG, Grant Thornton, Custom
- All previous features retained

NEW in v5.1:
- PDF & PPTX enabled by default in UI
- Speaker notes toggle: Minimal vs Detailed for PPTX
- All table/numbering/naming fixes retained
- Added Table of Contents toggle for both DOCX and PDF
- Smart TOC positioning at the BEGINNING of the document
- Smart heading detection for bold text lines
"""

import streamlit as st
import pypandoc
import tempfile
import os

from pathlib import Path
import re
from datetime import datetime
from docx import Document
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH, WD_COLOR_INDEX, WD_BREAK
from docx.enum.style import WD_STYLE_TYPE
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from pptx import Presentation
from pptx.util import Inches as PptInches, Pt as PptPt
from pptx.dml.color import RGBColor as PptRGBColor

# Ensure Pandoc/LaTeX see MacTeX binaries
os.environ["PATH"] = "/Library/TeX/texbin:" + os.environ.get("PATH", "")

# =========================
# GLOBAL CONFIG & COLOR THEMES
# =========================

GOVT_CONFIG = {
    "default_font": "Times New Roman",
    "body_size": 12,
    "heading_size": 14,
    "line_spacing": 1.5,
    "left_margin": 1.5,
    "navy_blue": RGBColor(0, 0, 128),
    "version": "5.2",
}

# Color theme presets
COLOR_THEMES = {
    "MoMSME (Navy Blue)": {
        "heading": (0, 0, 128),      # Navy Blue
        "table_header": "00008B",     # Dark Blue (hex for tables)
        "table_border": "4472C4",     # Medium Blue
        "description": "Government standard navy blue theme"
    },
    "KPMG (Blue)": {
        "heading": (0, 51, 141),      # KPMG Blue #00338D
        "table_header": "00338D",     # KPMG Blue
        "table_border": "00338D",     # KPMG Blue
        "description": "KPMG corporate brand colors"
    },
    "Grant Thornton (Purple)": {
        "heading": (80, 45, 127),     # GT Purple #502D7F
        "table_header": "502D7F",     # GT Purple
        "table_border": "502D7F",     # GT Purple
        "description": "Grant Thornton brand colors"
    },
    "Custom": {
        "heading": (0, 0, 128),       # Default, will be overridden
        "table_header": "4472C4",     # Default, will be overridden
        "table_border": "4472C4",     # Default
        "description": "Choose your own custom color"
    }
}

# =========================
# HELPERS – CLEANING / SLIDES
# =========================

def clean_ai_artifacts(raw_md: str) -> str:
    raw_md = re.sub(r"\[(web|cite):\d+\]", "", raw_md)
    raw_md = re.sub(r"\[\d+\](?!\s*https?://)", "", raw_md)

    emoji_pattern = re.compile(
        "["
        "\U0001F600-\U0001F64F"
        "\U0001F300-\U0001F5FF"
        "\U0001F680-\U0001F6FF"
        "\U0001F1E0-\U0001F1FF"
        "\U00002702-\U000027B0"
        "\U000024C2-\U0001F251"
        "]+",
        flags=re.UNICODE,
    )
    raw_md = emoji_pattern.sub("", raw_md)
    raw_md = raw_md.replace("—", "-").replace("–", "-").replace("―", "-")
    raw_md = re.sub(r"^ {2,10}", "", raw_md, flags=re.MULTILINE)
    raw_md = re.sub(r"\n{4,}", "\n\n\n", raw_md)
    raw_md = re.sub(r"```(?:python|text|bash|)\n?", "```\n", raw_md, flags=re.IGNORECASE)
    raw_md = "\n".join(line.rstrip() for line in raw_md.split("\n"))
    return raw_md.strip()


def prepare_slides_md(md_content: str, auto_breaks: bool = True, notes_style: str = "Minimal") -> str:
    """
    Prepare MD for slides with configurable speaker notes.

    notes_style:
        - "Minimal": Simple placeholder notes
        - "Detailed": Expanded talking points based on content
    """
    if auto_breaks:
        md_content = re.sub(r"^## ", "\n---\n\n## ", md_content, flags=re.MULTILINE)

    if notes_style == "Minimal":
        slide_pattern = r"(^##\s+.+?)(?=\n##|\n---|$)"

        def add_minimal_notes(match):
            slide = match.group(1)
            return f"{slide}\n\n::: notes\nKey talking points for this slide.\n:::"

        if ":::" not in md_content:
            md_content = re.sub(
                slide_pattern,
                add_minimal_notes,
                md_content,
                flags=re.MULTILINE | re.DOTALL,
            )

    elif notes_style == "Detailed":
        slide_pattern = r"(^##\s+(.+?))(\n(?:(?!^##).)*?)(?=\n##|\n---|$)"

        def add_detailed_notes(match):
            slide_heading = match.group(1)
            heading_text = match.group(2)
            slide_content = match.group(3) if match.group(3) else ""

            bullets = re.findall(r'^[-*]\s+(.+)$', slide_content, re.MULTILINE)

            notes = f"Speaker notes for: {heading_text}\n\n"
            if bullets:
                notes += "Key points to cover:\n"
                for bullet in bullets[:5]:
                    notes += f"- Elaborate on: {bullet}\n"
            else:
                notes += "Discuss the main concepts presented on this slide.\n"

            return f"{slide_heading}{slide_content}\n\n::: notes\n{notes}:::"

        if ":::" not in md_content:
            md_content = re.sub(
                slide_pattern,
                add_detailed_notes,
                md_content,
                flags=re.MULTILINE | re.DOTALL,
            )

    return md_content


# =========================
# PAGE NUMBERING & TOC
# =========================

def _add_field(run, instruction: str):
    fldChar1 = OxmlElement("w:fldChar")
    fldChar1.set(qn("w:fldCharType"), "begin")

    instrText = OxmlElement("w:instrText")
    instrText.set(qn("xml:space"), "preserve")
    instrText.text = instruction

    fldChar2 = OxmlElement("w:fldChar")
    fldChar2.set(qn("w:fldCharType"), "end")

    run._r.append(fldChar1)
    run._r.append(instrText)
    run._r.append(fldChar2)


def add_page_numbers(doc: Document, style: str = "1,2,3", position: str = "right"):
    section = doc.sections[0]
    footer = section.footer

    while footer.paragraphs:
        p = footer.paragraphs[0]._element
        footer._element.remove(p)

    footer_para = footer.add_paragraph()

    if position == "right":
        footer_para.alignment = WD_ALIGN_PARAGRAPH.RIGHT
    else:
        footer_para.alignment = WD_ALIGN_PARAGRAPH.CENTER

    if style == "Page X of Y":
        run = footer_para.add_run("Page ")
        _add_field(run, "PAGE")
        footer_para.add_run(" of ")
        run2 = footer_para.add_run()
        _add_field(run2, "NUMPAGES")

    elif style == "X of Y":
        run = footer_para.add_run()
        _add_field(run, "PAGE")
        footer_para.add_run(" of ")
        run2 = footer_para.add_run()
        _add_field(run2, "NUMPAGES")

    elif style == "1,2,3":
        run = footer_para.add_run()
        _add_field(run, "PAGE")

    for run in footer_para.runs:
        run.font.name = "Times New Roman"
        run.font.size = Pt(10)
        run.font.color.rgb = RGBColor(0, 0, 0)


def insert_docx_toc(doc: Document, heading_color: tuple):
    """
    Inserts a Table of Contents field at the BEGINNING of the document.
    """
    if not doc.paragraphs:
        return
        
    # Get the very first paragraph in the document
    first_para = doc.paragraphs[0]
    
    # 1. Insert formatted title BEFORE the first paragraph
    toc_heading = first_para.insert_paragraph_before()
    toc_heading.alignment = WD_ALIGN_PARAGRAPH.LEFT
    run_heading = toc_heading.add_run("Table of Contents")
    run_heading.font.name = "Times New Roman"
    run_heading.font.size = Pt(16)
    run_heading.font.bold = True
    run_heading.font.color.rgb = RGBColor(*heading_color)
    
    # 2. Add TOC field code BEFORE the first paragraph (after the title)
    toc_para = first_para.insert_paragraph_before()
    run = toc_para.add_run()

    fldChar1 = OxmlElement('w:fldChar')
    fldChar1.set(qn('w:fldCharType'), 'begin')

    instrText = OxmlElement('w:instrText')
    instrText.set(qn('xml:space'), 'preserve')
    instrText.text = 'TOC \\o "1-3" \\h \\z \\u'

    fldChar2 = OxmlElement('w:fldChar')
    fldChar2.set(qn('w:fldCharType'), 'separate')

    fldChar3 = OxmlElement('w:fldChar')
    fldChar3.set(qn('w:fldCharType'), 'end')

    run._r.append(fldChar1)
    run._r.append(instrText)
    run._r.append(fldChar2)
    run._r.append(fldChar3)
    
    # 3. Add page break BEFORE the first paragraph
    pb_para = first_para.insert_paragraph_before()
    pb_para.add_run().add_break(WD_BREAK.PAGE)


# =========================
# TABLE STYLING (FIX WHITE TEXT)
# =========================

def fix_table_font_colors_preserve_background(doc: Document):
    for table in doc.tables:
        for row_idx, row in enumerate(table.rows):
            for cell in row.cells:
                tcPr = cell._element.get_or_add_tcPr()
                shd_elems = tcPr.findall(qn("w:shd"))
                cell_has_bg = bool(shd_elems)

                for para in cell.paragraphs:
                    for run in para.runs:
                        color = run.font.color.rgb if run.font.color is not None else None

                        if row_idx == 0:
                            if color is None:
                                run.font.color.rgb = RGBColor(255, 255, 255)
                            continue

                        if color is None or (color[0], color[1], color[2]) == (255, 255, 255):
                            run.font.color.rgb = RGBColor(0, 0, 0)


def apply_table_look_and_feel(doc: Document, table_header_color: str, table_border_color: str):
    """
    Apply themed table styling with custom colors.
    
    Args:
        table_header_color: Hex color for header background (e.g., "4472C4")
        table_border_color: Hex color for borders (e.g., "4472C4")
    """
    for table in doc.tables:
        tbl = table._tbl
        tblPr = tbl.tblPr

        tblBorders = OxmlElement("w:tblBorders")
        for border_name in ["top", "left", "bottom", "right", "insideH", "insideV"]:
            border = OxmlElement(f"w:{border_name}")
            border.set(qn("w:val"), "single")
            border.set(qn("w:sz"), "4")
            border.set(qn("w:space"), "0")
            border.set(qn("w:color"), table_border_color)
            tblBorders.append(border)
        tblPr.append(tblBorders)

        # Header row
        if len(table.rows) > 0:
            hdr_cells = table.rows[0].cells
            for cell in hdr_cells:
                tcPr = cell._element.get_or_add_tcPr()
                shd_elems = tcPr.findall(qn("w:shd"))
                if not shd_elems:
                    shading_elm = OxmlElement("w:shd")
                    shading_elm.set(qn("w:fill"), table_header_color)
                    tcPr.append(shading_elm)

                for paragraph in cell.paragraphs:
                    for run in paragraph.runs:
                        run.font.bold = True
                        run.font.color.rgb = RGBColor(255, 255, 255)
                        if run.font.name is None:
                            run.font.name = "Times New Roman"
                        if run.font.size is None:
                            run.font.size = Pt(11)

        # Body rows
        for row_idx, row in enumerate(table.rows[1:], start=1):
            for cell in row.cells:
                for paragraph in cell.paragraphs:
                    for run in paragraph.runs:
                        if run.font.name is None:
                            run.font.name = "Times New Roman"
                        if run.font.size is None:
                            run.font.size = Pt(11)

                tcPr = cell._element.get_or_add_tcPr()
                shd_elems = tcPr.findall(qn("w:shd"))
                if row_idx % 2 == 0 and not shd_elems:
                    shading_elm = OxmlElement("w:shd")
                    shading_elm.set(qn("w:fill"), "D9E2F3")
                    tcPr.append(shading_elm)


# =========================
# HEADING RESTYLE FOR UPLOADED DOCX
# =========================

def restyle_docx_headings(doc: Document, heading_color=(0, 0, 128)):
    """
    Make headings use theme color TNR.
    - True Word headings: style name starts with 'Heading'
    - Pseudo-headings: Detects standalone lines under 150 chars where all text is bold.
    """
    theme_color = RGBColor(*heading_color)
    
    for para in doc.paragraphs:
        txt = para.text.strip()
        style = para.style
        sname = style.name.lower() if style is not None and style.name else ""
        
        is_word_heading = sname.startswith("heading")
        
        # Detect bold standalone lines (pseudo-headings)
        is_pseudo_heading = False
        if txt and len(txt) < 150 and not is_word_heading:
            # Check if all runs with actual text are bold
            text_runs = [r for r in para.runs if r.text.strip()]
            if text_runs and all(r.bold for r in text_runs):
                is_pseudo_heading = True
        
        if not (is_word_heading or is_pseudo_heading):
            continue
        
        # basic level‑based sizing
        if "1" in sname:
            size = Pt(16)
        elif "2" in sname:
            size = Pt(14)
        elif is_pseudo_heading:
            size = Pt(14)  # Make bold pseudo-headings 14pt by default
        else:
            size = Pt(13)
            
        para.alignment = WD_ALIGN_PARAGRAPH.LEFT
        for run in para.runs:
            run.font.name = "Times New Roman"
            run.font.size = size
            run.font.bold = True
            run.font.color.rgb = theme_color
            # Clear any highlight to keep it clean
            run.font.highlight_color = None
            

# =========================
# METADATA & TEMPLATE DOCX
# =========================

def set_document_properties(doc: Document, author: str = "Sandeep Prasad", title: str | None = None):
    core_props = doc.core_properties
    core_props.author = author
    
    # Truncate title to avoid the 255-char limit crash in Word properties
    safe_title = title or "MoMSME Document"
    if len(safe_title) > 250:
        safe_title = safe_title[:247] + "..."
        
    core_props.title = safe_title
    core_props.subject = "Official Government Document"
    core_props.keywords = "MSME, Government of India"
    now = datetime.now()
    core_props.created = now
    core_props.modified = now


def create_docx_template(
    path: str,
    header_text: str,
    font_body: str,
    font_heading: str,
    body_size: int,
    heading_size: int,
    body_color,
    heading_color,
    align: str,
    bold_body: bool,
    italic_body: bool,
    underline_body: bool,
    page_number_style: str,
    author: str,
) -> str:
    doc = Document()
    set_document_properties(doc, author=author)

    section = doc.sections[0]
    section.left_margin = Inches(GOVT_CONFIG["left_margin"])
    section.right_margin = Inches(1.0)
    section.top_margin = Inches(1.0)
    section.bottom_margin = Inches(1.0)

    body_align = WD_ALIGN_PARAGRAPH.JUSTIFY if align == "Justify" else WD_ALIGN_PARAGRAPH.LEFT
    body_rgb = RGBColor(*body_color)
    heading_rgb = RGBColor(*heading_color)

    styles = doc.styles

    body_style = styles.add_style("GovtBody", WD_STYLE_TYPE.PARAGRAPH)
    body_style.font.name = font_body
    body_style.font.size = Pt(body_size)
    body_style.font.bold = bold_body
    body_style.font.italic = italic_body
    body_style.font.underline = underline_body
    body_style.font.color.rgb = body_rgb
    body_style.paragraph_format.line_spacing = GOVT_CONFIG["line_spacing"]
    body_style.paragraph_format.alignment = body_align

    for level in range(1, 4):
        h_style = styles.add_style(f"GovtHeading{level}", WD_STYLE_TYPE.PARAGRAPH)
        h_style.font.name = font_heading
        h_style.font.size = Pt(heading_size + (3 - level) * 2)
        h_style.font.bold = True
        h_style.font.color.rgb = heading_rgb
        h_style.paragraph_format.space_after = Pt(12)
        h_style.paragraph_format.alignment = WD_ALIGN_PARAGRAPH.LEFT

    header = doc.add_paragraph("GOVERNMENT OF INDIA", style="GovtHeading1")
    header.alignment = WD_ALIGN_PARAGRAPH.CENTER

    sub = doc.add_paragraph(header_text, style="GovtBody")
    sub.alignment = WD_ALIGN_PARAGRAPH.CENTER

    subsub = doc.add_paragraph("IT & Planning Division", style="GovtBody")
    subsub.alignment = WD_ALIGN_PARAGRAPH.CENTER

    doc.add_page_break()

    if page_number_style == "1, 2, 3, ...":
        add_page_numbers(doc, style="1,2,3", position="right")
    elif page_number_style == "Page X of Y":
        add_page_numbers(doc, style="Page X of Y", position="right")
    elif page_number_style == "X of Y":
        add_page_numbers(doc, style="X of Y", position="right")

    doc.save(path)
    return path


def create_pptx_template(path: str, heading_color: tuple) -> str:
    prs = Presentation()
    prs.slide_width = PptInches(10)
    prs.slide_height = PptInches(7.5)

    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)

    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Ministry of MSME"
    title_p = title.text_frame.paragraphs[0]
    title_p.font.name = GOVT_CONFIG["default_font"]
    title_p.font.size = PptPt(44)
    title_p.font.color.rgb = PptRGBColor(*heading_color)

    subtitle.text = "IT & Planning Division"
    subtitle_p = subtitle.text_frame.paragraphs[0]
    subtitle_p.font.name = GOVT_CONFIG["default_font"]
    subtitle_p.font.size = PptPt(28)
    subtitle_p.font.color.rgb = PptRGBColor(0, 0, 0)

    prs.save(path)
    return path


# =========================
# JUSTIFY + HIGHLIGHT
# =========================

def force_justify(doc: Document):
    for para in doc.paragraphs:
        if para.text.strip():
            para.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY


def apply_yellow_highlight_from_markers(doc: Document, marker_pattern: str = r"==(.+?)=="):
    regex = re.compile(marker_pattern)

    for para in doc.paragraphs:
        text = para.text
        if "==" not in text:
            continue

        matches = list(regex.finditer(text))
        if not matches:
            continue

        new_runs = []
        current_pos = 0

        for match in matches:
            start, end = match.span()
            before = text[current_pos:start]
            highlighted = match.group(1)
            current_pos = end

            if before:
                new_runs.append(("normal", before))
            if highlighted:
                new_runs.append(("highlight", highlighted))

        after = text[current_pos:]
        if after:
            new_runs.append(("normal", after))

        for r in list(para.runs):
            r.clear()
        para.text = ""

        for kind, segment in new_runs:
            run = para.add_run(segment)
            if kind == "highlight":
                run.font.highlight_color = WD_COLOR_INDEX.YELLOW


# =========================
# TITLE + TIMESTAMP + FILENAMES
# =========================

def extract_title_from_md(md: str) -> str | None:
    lines = [l.strip() for l in md.splitlines() if l.strip()]
    for line in lines:
        if line.startswith("# "):
            return line[2:].strip()
    for line in lines:
        if len(line) > 20:
            return line.strip()
    return None


def extract_title_from_docx(doc: Document) -> str | None:
    for para in doc.paragraphs[:25]:
        txt = para.text.strip()
        if txt and len(txt) > 20:
            return txt
    return None


def sanitize_filename_part(text: str) -> str:
    text = re.sub(r"\s+", " ", text).strip()
    text = text.replace("/", "-").replace("\\", "-")
    text = re.sub(r"[^A-Za-z0-9\-\s]", "", text)
    text = text[:80]
    text = text.replace(" ", "-")
    return text or "MoMSME-Document"


def get_timestamp() -> str:
    return datetime.now().strftime("%Y%m%d_%H%M")


def build_output_name(base_override: str | None, auto_title: str | None, suffix: str) -> str:
    if base_override and base_override.strip():
        base = sanitize_filename_part(base_override)
    else:
        base = sanitize_filename_part(auto_title or "MoMSME-Document")
    ts = get_timestamp()
    return f"{base}_{ts}.{suffix}"


# =========================
# GENERATION – MD SOURCE
# =========================

def generate_documents_from_md(content: str, output_formats, options, filename_base: str | None):
    results = {}

    with tempfile.TemporaryDirectory() as tmp:
        tmp_path = Path(tmp)
        auto_title = extract_title_from_md(content)

        if "DOCX" in output_formats:
            docx_template = tmp_path / "template.docx"
            docx_output = tmp_path / "output.docx"

            create_docx_template(
                path=str(docx_template),
                header_text=options["header"],
                font_body=options["font_body"],
                font_heading=options["font_heading"],
                body_size=options["body_size"],
                heading_size=options["heading_size"],
                body_color=options["body_color"],
                heading_color=options["heading_color"],
                align=options["align"],
                bold_body=options["bold_body"],
                italic_body=options["italic_body"],
                underline_body=options["underline_body"],
                page_number_style=options["page_number_style"],
                author=options["author"],
            )

            extra_args = [
                f"--reference-doc={docx_template}",
                "--from=markdown+pipe_tables+simple_tables+grid_tables+citations",
            ]
            if options["toc"]:
                extra_args.append("--toc")

            pypandoc.convert_text(
                content,
                to="docx",
                format="md",
                outputfile=str(docx_output),
                extra_args=extra_args,
            )

            doc = Document(str(docx_output))
            force_justify(doc)
            apply_table_look_and_feel(
                doc, 
                table_header_color=options["table_header_color"],
                table_border_color=options["table_border_color"]
            )
            fix_table_font_colors_preserve_background(doc)
            set_document_properties(doc, author=options["author"], title=auto_title)
            apply_yellow_highlight_from_markers(doc)
            doc.save(str(docx_output))

            with open(docx_output, "rb") as f:
                fname = build_output_name(filename_base, auto_title, "docx")
                results["DOCX"] = (fname, f.read())

        if "PPTX" in output_formats:
            pptx_template = tmp_path / "ppt_template.pptx"
            pptx_output = tmp_path / "output.pptx"

            if options["govt_template"]:
                create_pptx_template(str(pptx_template), options["heading_color"])

            slide_content = prepare_slides_md(
                content,
                options["auto_breaks"],
                notes_style=options.get("pptx_notes_style", "Minimal"),
            )

            extra_args = [
                "--slide-level=2",
                "--from=markdown+pipe_tables",
            ]
            if options["govt_template"]:
                extra_args.append(f"--reference-doc={pptx_template}")

            pypandoc.convert_text(
                slide_content,
                to="pptx",
                format="md",
                outputfile=str(pptx_output),
                extra_args=extra_args,
            )

            with open(pptx_output, "rb") as f:
                fname = build_output_name(filename_base, auto_title, "pptx")
                results["PPTX"] = (fname, f.read())

        if "PDF" in output_formats:
            pdf_output = tmp_path / "output.pdf"

            extra_args = [
                "--pdf-engine=xelatex",
                f"--variable=mainfont={options['font_body']}",
                f"--variable=fontsize={options['body_size']}pt",
                "--variable=geometry:margin=1.5in",
            ]
            
            # Use Pandoc's built-in TOC for PDF
            if options["toc"]:
                extra_args.append("--toc")

            pypandoc.convert_text(
                content,
                to="pdf",
                format="md",
                outputfile=str(pdf_output),
                extra_args=extra_args,
            )

            with open(pdf_output, "rb") as f:
                fname = build_output_name(filename_base, auto_title, "pdf")
                results["PDF"] = (fname, f.read())

    return results


# =========================
# GENERATION – DOCX SOURCE
# =========================

def generate_documents_from_docx(source_docx_bytes: bytes, output_formats, options, filename_base: str | None):
    results = {}

    with tempfile.TemporaryDirectory() as tmp:
        tmp_path = Path(tmp)
        src_path = tmp_path / "input.docx"
        with open(src_path, "wb") as f:
            f.write(source_docx_bytes)

        doc = Document(str(src_path))
        auto_title = extract_title_from_docx(doc)

        restyle_docx_headings(doc, heading_color=options["heading_color"])

        force_justify(doc)
        apply_table_look_and_feel(
            doc,
            table_header_color=options["table_header_color"],
            table_border_color=options["table_border_color"]
        )
        fix_table_font_colors_preserve_background(doc)
        apply_yellow_highlight_from_markers(doc)

        if options["page_number_style"] != "None":
            style_map = {
                "1, 2, 3, ...": "1,2,3",
                "Page X of Y": "Page X of Y",
                "X of Y": "X of Y",
            }
            style = style_map.get(options["page_number_style"], "1,2,3")
            add_page_numbers(doc, style=style, position="right")
            
        # Insert TOC at the TOP of the doc if requested
        if options["toc"]:
            insert_docx_toc(doc, heading_color=options["heading_color"])

        set_document_properties(doc, author=options["author"], title=auto_title)
        enhanced_path = tmp_path / "enhanced.docx"
        doc.save(str(enhanced_path))

        if "DOCX" in output_formats:
            with open(enhanced_path, "rb") as f:
                fname = build_output_name(filename_base, auto_title, "docx")
                results["DOCX"] = (fname, f.read())

        if "PPTX" in output_formats or "PDF" in output_formats:
            temp_md = tmp_path / "temp.md"
            pypandoc.convert_file(
                str(enhanced_path),
                to="markdown",
                format="docx",
                outputfile=str(temp_md),
            )
            with open(temp_md, "r", encoding="utf-8") as f_md:
                md_content = clean_ai_artifacts(f_md.read())

            md_results = generate_documents_from_md(md_content, output_formats, options, filename_base)

            for fmt in ["PPTX", "PDF"]:
                if fmt in md_results:
                    results[fmt] = md_results[fmt]

    return results


# =========================
# STREAMLIT APP
# =========================

def main():
    st.title("🎨 MoMSME Document Formatter v5.2")
    st.caption(
        "Full-featured formatter with MoMSME template, color themes (KPMG, Grant Thornton, Custom), "
        "tables/images preserved, smart filenames, PPTX speaker notes, and Table of Contents."
    )

    mode = st.radio("Input Type", ["Markdown", "DOCX"])

    output_formats = st.multiselect(
        "Output Formats",
        ["DOCX", "PPTX", "PDF"],
        default=["DOCX", "PPTX", "PDF"],
    )

    st.markdown("---")
    st.subheader("🎨 Color Theme")
    
    # Color theme selector
    theme_choice = st.selectbox(
        "Select Color Theme",
        list(COLOR_THEMES.keys()),
        index=0,
        help="Choose a predefined color theme or select 'Custom' to pick your own colors"
    )
    
    # Show theme description
    st.caption(f"_{COLOR_THEMES[theme_choice]['description']}_")
    
    # If Custom is selected, show color picker
    if theme_choice == "Custom":
        st.markdown("**Custom Color Selection**")
        col_a, col_b = st.columns(2)
        with col_a:
            custom_heading = st.color_picker(
                "Heading Color",
                value="#00008B",
                help="Choose the color for headings and titles"
            )
        with col_b:
            custom_table = st.color_picker(
                "Table Header Color",
                value="#4472C4",
                help="Choose the color for table headers and borders"
            )
        
        # Convert hex to RGB tuple
        heading_rgb = tuple(int(custom_heading.lstrip('#')[i:i+2], 16) for i in (0, 2, 4))
        table_hex = custom_table.lstrip('#').upper()
        
        selected_theme = {
            "heading": heading_rgb,
            "table_header": table_hex,
            "table_border": table_hex,
        }
    else:
        selected_theme = COLOR_THEMES[theme_choice]
    
    # Show color preview
    st.markdown("**Theme Preview:**")
    preview_col1, preview_col2, preview_col3 = st.columns(3)
    with preview_col1:
        rgb_str = f"rgb{selected_theme['heading']}"
        st.markdown(f"<div style='background-color: {rgb_str}; padding: 20px; border-radius: 5px; color: white; text-align: center;'><strong>Heading Color</strong></div>", unsafe_allow_html=True)
    with preview_col2:
        st.markdown(f"<div style='background-color: #{selected_theme['table_header']}; padding: 20px; border-radius: 5px; color: white; text-align: center;'><strong>Table Header</strong></div>", unsafe_allow_html=True)
    with preview_col3:
        st.markdown(f"<div style='background-color: #{selected_theme['table_border']}; padding: 20px; border-radius: 5px; color: white; text-align: center;'><strong>Table Border</strong></div>", unsafe_allow_html=True)

    st.markdown("---")
    st.subheader("📝 Look & Feel Settings")

    col1, col2 = st.columns(2)
    with col1:
        header_text = st.text_input("Header Ministry Text", "Ministry of MSME")
        font_body = st.selectbox("Body Font", ["Times New Roman", "Calibri", "Arial"], index=0)
        font_heading = st.selectbox("Heading Font", ["Times New Roman", "Calibri", "Arial"], index=0)
        body_size = st.slider("Body Font Size", 10, 14, GOVT_CONFIG["body_size"])
        heading_size = st.slider("Heading Base Size", 12, 18, GOVT_CONFIG["heading_size"])

    with col2:
        align = st.selectbox("Body Alignment", ["Justify", "Left"], index=0)
        bold_body = st.checkbox("Body Bold", False)
        italic_body = st.checkbox("Body Italic", False)
        underline_body = st.checkbox("Body Underline", False)
        page_number_style = st.selectbox(
            "Page Numbering Style",
            ["None", "1, 2, 3, ...", "Page X of Y", "X of Y"],
            index=2,
        )
        include_toc = st.checkbox("Include Table of Contents", False)

    st.subheader("📊 PPTX Settings")
    pptx_notes_style = st.radio(
        "Speaker Notes Style",
        ["Minimal", "Detailed"],
        index=0,
        help="Minimal: Simple placeholder notes. Detailed: Expanded talking points based on content.",
    )

    st.subheader("💾 Output Filename")
    filename_base = st.text_input(
        "Base name for output files (optional):",
        "",
        help="If empty, the app will infer a title from the document and append a timestamp.",
    )

    options = {
        "header": header_text,
        "font_body": font_body,
        "font_heading": font_heading,
        "body_size": body_size,
        "heading_size": heading_size,
        "body_color": (0, 0, 0),
        "heading_color": selected_theme["heading"],
        "table_header_color": selected_theme["table_header"],
        "table_border_color": selected_theme["table_border"],
        "align": align,
        "bold_body": bold_body,
        "italic_body": italic_body,
        "underline_body": underline_body,
        "page_number_style": page_number_style,
        "author": "Sandeep Prasad",
        "toc": include_toc,
        "beautify_tables": True,
        "govt_template": True,
        "auto_breaks": True,
        "pptx_notes_style": pptx_notes_style,
    }

    st.markdown("---")

    if mode == "Markdown":
        md_input = st.text_area("Paste Markdown Content", height=400)
        if st.button("🚀 Generate from Markdown", type="primary") and md_input.strip():
            with st.spinner("Generating documents..."):
                cleaned = clean_ai_artifacts(md_input)
                results = generate_documents_from_md(cleaned, output_formats, options, filename_base)
            
            st.success("✅ Documents generated successfully!")
            for fmt, (fname, blob) in results.items():
                st.download_button(
                    label=f"📥 Download {fmt} – {fname}",
                    data=blob,
                    file_name=fname,
                    mime=(
                        "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
                        if fmt == "DOCX"
                        else "application/vnd.openxmlformats-officedocument.presentationml.presentation"
                        if fmt == "PPTX"
                        else "application/pdf"
                    ),
                )
    else:
        uploaded = st.file_uploader("Upload DOCX", type=["docx"])
        if uploaded is not None and st.button("🚀 Generate from DOCX", type="primary"):
            with st.spinner("Processing document..."):
                src_bytes = uploaded.read()
                results = generate_documents_from_docx(src_bytes, output_formats, options, filename_base)
            
            st.success("✅ Documents generated successfully!")
            for fmt, (fname, blob) in results.items():
                st.download_button(
                    label=f"📥 Download {fmt} – {fname}",
                    data=blob,
                    file_name=fname,
                    mime=(
                        "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
                        if fmt == "DOCX"
                        else "application/vnd.openxmlformats-officedocument.presentationml.presentation"
                        if fmt == "PPTX"
                        else "application/pdf"
                    ),
                )

if __name__ == "__main__":
    main()
